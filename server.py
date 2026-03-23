import os
import re
import io
import csv
import base64
import json
import traceback
import requests as http_req
from flask import Flask, request, jsonify, send_from_directory, session, g, has_request_context
from dotenv import load_dotenv
from markdownify import markdownify as md_convert

load_dotenv()

app = Flask(__name__, static_folder='public', static_url_path='')
app.config['MAX_CONTENT_LENGTH'] = int(os.getenv('MAX_CONTENT_LENGTH', str(25 * 1024 * 1024)))
app.secret_key = (os.getenv('FLASK_SECRET_KEY') or '').strip() or os.urandom(24).hex()

CONFLUENCE_URL = (os.getenv('CONFLUENCE_URL') or '').rstrip('/')
DEFAULT_EMAIL_DOMAIN = (os.getenv('DEFAULT_EMAIL_DOMAIN') or 'vuno.co').strip()

api_base_path = '/rest/api'

# 채팅에 허용하는 OpenAI 모델 (이 3개만 선택·호출 가능)
ALLOWED_OPENAI_MODELS = ('gpt-5.4', 'gpt-5.4-mini', 'gpt-5.4-nano')
ALLOWED_OPENAI_MODELS_SET = frozenset(ALLOWED_OPENAI_MODELS)

# 첨부 파일 (채팅)
CHAT_ATTACHMENT_MAX_FILES = int(os.getenv('CHAT_ATTACHMENT_MAX_FILES', '8'))
CHAT_ATTACHMENT_MAX_BYTES = int(os.getenv('CHAT_ATTACHMENT_MAX_BYTES', str(20 * 1024 * 1024)))  # per file
CHAT_ATTACHMENT_MAX_EXTRACT_CHARS = int(os.getenv('CHAT_ATTACHMENT_MAX_EXTRACT_CHARS', str(100_000)))  # per file text

IMAGE_MIME_PREFIX = 'image/'
VISION_IMAGE_TYPES = frozenset({
    'image/jpeg', 'image/png', 'image/gif', 'image/webp',
})

_env_model = (os.getenv('AI_MODEL') or 'gpt-5.4-mini').strip()
if _env_model not in ALLOWED_OPENAI_MODELS_SET:
    _env_model = 'gpt-5.4-mini'

ai_config = {
    'provider': 'openai',
    'api_key': '',
    'base_url': '',
    'model': _env_model,
    'temperature': 0.2,
    'top_p': 0.9,
    'repetition_penalty': 1.0,
    'max_tokens': 4096,
}


def _normalize_confluence_login_id(raw):
    raw = (raw or '').strip()
    if not raw:
        return ''
    if '@' in raw:
        return raw
    if DEFAULT_EMAIL_DOMAIN:
        return f'{raw}@{DEFAULT_EMAIL_DOMAIN}'
    return raw


@app.before_request
def _bind_request_credentials():
    g.openai_api_key = (session.get('openai_api_key') or '').strip()
    g.cf_token = None
    g.cf_username = ''
    g.cf_auth_mode = None
    g.cf_api_base = None
    if session.get('cf_token'):
        g.cf_token = session['cf_token']
        g.cf_username = (session.get('cf_username') or '').strip()
        g.cf_auth_mode = (session.get('cf_auth_mode') or 'basic').strip().lower()
        g.cf_api_base = (session.get('cf_api_base') or '').strip() or None


def get_openai_key():
    if has_request_context():
        v = (getattr(g, 'openai_api_key', None) or '').strip()
        if v:
            return v
    return (ai_config.get('api_key') or '').strip()


def get_effective_api_base_path():
    if has_request_context():
        b = getattr(g, 'cf_api_base', None)
        if b:
            return b
    return api_base_path


def get_auth():
    if has_request_context():
        tok = getattr(g, 'cf_token', None)
        mode = getattr(g, 'cf_auth_mode', None)
        if tok and mode == 'basic':
            u = (getattr(g, 'cf_username', None) or '').strip()
            if u:
                return (u, tok)
    return None


def get_headers():
    headers = {
        'Content-Type': 'application/json',
        'Accept': 'application/json',
        'X-Atlassian-Token': 'no-check',
    }
    if has_request_context():
        tok = getattr(g, 'cf_token', None)
        mode = getattr(g, 'cf_auth_mode', None)
        if tok and mode == 'bearer':
            headers['Authorization'] = f'Bearer {tok}'
    return headers


def is_configured():
    if not CONFLUENCE_URL:
        return False
    if has_request_context() and getattr(g, 'cf_token', None):
        return True
    return False


def _confluence_json_dict(resp):
    if resp is None or resp.status_code != 200:
        return None
    ct = (resp.headers.get('Content-Type') or '').lower()
    if 'json' not in ct:
        return None
    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError):
        return None
    return data if isinstance(data, dict) else None


def _is_confluence_anonymous_user(data):
    if not data:
        return True
    t = (data.get('type') or '').lower()
    if t == 'anonymous':
        return True
    u = (data.get('username') or '').lower()
    if u == 'anonymous':
        return True
    at = (data.get('accountType') or '').lower()
    if at == 'anonymous':
        return True
    return False


def _greeting_token_from_display_name(display_name):
    """
    프로필 표시명이 '한국어 3자 / 영어 이름' 형태면 슬래시 앞 한국어 부분의 뒤 2글자만 사용.
    (예: 이나혁 / Nahyuk Lee -> 나혁)
    """
    s = (display_name or '').strip()
    if not s or '/' not in s:
        return None
    korean_part = s.split('/', 1)[0].strip()
    if not korean_part:
        return None
    if len(korean_part) >= 2:
        return korean_part[-2:]
    return korean_part


def _confluence_greeting_name(user_json, raw_login_id, expected_email):
    """웰컴 인사에 붙일 짧은 호칭 (한/영 표시명이면 한국어 뒤 2글자, 아니면 영문 첫 덩어리·ID 등)."""
    if isinstance(user_json, dict):
        for key in ('displayName', 'publicName'):
            v = user_json.get(key)
            if isinstance(v, str) and v.strip():
                vv = v.strip()
                tok = _greeting_token_from_display_name(vv)
                if tok:
                    return tok
                first = vv.split()[0] if vv.split() else vv
                if first:
                    return first
    rid = (raw_login_id or '').strip()
    if rid:
        return rid.split('@', 1)[0].strip() if '@' in rid else rid
    em = (expected_email or '').strip()
    if '@' in em:
        return em.split('@', 1)[0].strip()
    return em or '사용자'


def _confluence_identity_proves_login(j_anon, j_auth, r_anon, r_auth, expected_email_l):
    """
    자격 증명이 실제로 반영됐는지 검증합니다.
    - 공개 /space 처럼 토큰 없이도 되는 엔드포인트와 달리 /user/current 기준.
    - 익명과 동일한 본문이면(Authorization 무시) 실패.
    - API가 email을 주면 입력한 ID(이메일)와 반드시 일치.
    """
    if not j_auth or _is_confluence_anonymous_user(j_auth):
        return False

    em = (j_auth.get('email') or '').strip().lower()
    if em:
        return em == expected_email_l

    if r_anon is not None and r_auth is not None:
        if r_anon.status_code == 200 and r_auth.status_code == 200:
            if (r_anon.text or '') == (r_auth.text or ''):
                return False

    if r_anon is not None and r_anon.status_code == 401:
        return True

    if j_anon is None or _is_confluence_anonymous_user(j_anon):
        return True

    aid_n = j_anon.get('accountId') or j_anon.get('userKey') or j_anon.get('username')
    aid_a = j_auth.get('accountId') or j_auth.get('userKey') or j_auth.get('username')
    if aid_n and aid_a and aid_n != aid_a:
        return True
    if aid_n and aid_a and aid_n == aid_a:
        return False
    return True


def probe_confluence_login(expected_email, token, raw_login_id=''):
    """
    /rest/api/user/current (및 wiki 베이스)로 실제 로그인 여부를 검사합니다.
    /space 는 익명 공개인 경우가 많아 비밀번호가 틀려도 200이 될 수 있어 사용하지 않습니다.
    성공 시 (api_base, auth_mode, greeting_name) 반환.
    """
    std_headers = {
        'Content-Type': 'application/json',
        'Accept': 'application/json',
        'X-Atlassian-Token': 'no-check',
    }
    last_err = None
    exp = (expected_email or '').strip().lower()
    if not exp:
        raise RuntimeError('사용자 ID(이메일)가 비어 있습니다.')

    for base in ('/rest/api', '/wiki/rest/api'):
        url = f'{CONFLUENCE_URL}{base}/user/current'

        try:
            r_anon = http_req.get(url, headers=std_headers, timeout=15)
        except Exception as e:
            r_anon = None
            last_err = str(e)
        j_anon = _confluence_json_dict(r_anon) if r_anon and r_anon.ok else None

        try:
            r_basic = http_req.get(
                url,
                headers=std_headers,
                auth=(expected_email, token),
                timeout=15,
            )
        except Exception as e:
            last_err = str(e)
            r_basic = None

        if r_basic is not None:
            if r_basic.status_code == 401:
                last_err = 'Confluence ID 또는 API 토큰이 올바르지 않습니다(401).'
            else:
                j_b = _confluence_json_dict(r_basic)
                if j_b and _confluence_identity_proves_login(j_anon, j_b, r_anon, r_basic, exp):
                    gname = _confluence_greeting_name(j_b, raw_login_id, expected_email)
                    return base, 'basic', gname
                if j_b and _is_confluence_anonymous_user(j_b):
                    last_err = (
                        '토큰이 인식되지 않아 익명으로만 접근됩니다. '
                        'Atlassian API 토큰(비밀번호 자리)과 이메일(ID) 조합을 확인하세요.'
                    )
                elif r_basic.ok:
                    last_err = (
                        '로그인된 계정이 입력한 이메일과 일치하지 않거나, '
                        '자격 증명이 API에 반영되지 않은 것 같습니다.'
                    )
                else:
                    last_err = f'user/current HTTP {r_basic.status_code}: {(r_basic.text or "")[:400]}'

        try:
            h = {**std_headers, 'Authorization': f'Bearer {token}'}
            r_bearer = http_req.get(url, headers=h, timeout=15)
        except Exception as e:
            last_err = str(e)
            continue

        if r_bearer.status_code == 401:
            last_err = 'Confluence Bearer 토큰이 올바르지 않습니다(401).'
            continue

        j_br = _confluence_json_dict(r_bearer)
        if j_br and _confluence_identity_proves_login(j_anon, j_br, r_anon, r_bearer, exp):
            gname = _confluence_greeting_name(j_br, raw_login_id, expected_email)
            return base, 'bearer', gname
        if j_br and _is_confluence_anonymous_user(j_br):
            last_err = 'Bearer 토큰이 유효하지 않거나 익명 응답만 받았습니다.'
        elif r_bearer.ok:
            last_err = 'Pat/Bearer로 로그인한 계정 이메일이 입력한 ID와 일치하지 않습니다.'

    raise RuntimeError(last_err or 'Confluence API에 연결할 수 없습니다.')


def verify_openai_key_for_login(api_key):
    """
    OpenAI 키로 /v1/models 를 호출하고, 허용 모델 ID 중 최소 하나가 계정에 있어야 통과합니다.
    """
    resp = http_req.get(
        'https://api.openai.com/v1/models',
        headers={'Authorization': f'Bearer {api_key}'},
        timeout=20,
    )
    if resp.status_code == 401:
        raise RuntimeError('OpenAI API 키가 거부되었습니다(401). 키를 확인하세요.')
    if not resp.ok:
        snippet = (resp.text or '')[:400]
        raise RuntimeError(f'OpenAI API 오류 (HTTP {resp.status_code}): {snippet or "본문 없음"}')
    ct = (resp.headers.get('Content-Type') or '').lower()
    if 'json' not in ct:
        raise RuntimeError('OpenAI 응답이 JSON이 아닙니다. 프록시·차단 여부를 확인하세요.')
    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError) as e:
        raise RuntimeError(f'OpenAI models JSON 파싱 실패: {e}') from e
    items = data.get('data')
    if not isinstance(items, list):
        raise RuntimeError('OpenAI /v1/models 응답에 data 배열이 없습니다.')
    api_ids = {m.get('id', '') for m in items if isinstance(m, dict)}
    allowed_hits = [mid for mid in ALLOWED_OPENAI_MODELS if mid in api_ids]
    if not allowed_hits:
        raise RuntimeError(
            f'이 키로는 앱에서 허용한 모델 중 하나도 쓸 수 없습니다. '
            f'필요: {", ".join(ALLOWED_OPENAI_MODELS)}. '
            '조직 정책·프로젝트 제한 또는 다른 키인지 확인하세요.'
        )
    return allowed_hits


def confluence_api(method, endpoint, body=None):
    base = get_effective_api_base_path()
    url = f'{CONFLUENCE_URL}{base}{endpoint}'
    kwargs = {'headers': get_headers(), 'timeout': 30}
    auth = get_auth()
    if auth:
        kwargs['auth'] = auth
    if body:
        kwargs['json'] = body
    resp = http_req.request(method, url, **kwargs)
    resp.raise_for_status()
    if resp.status_code == 204:
        return None
    return resp.json()


def html_to_markdown(html):
    if not html:
        return ''
    return md_convert(html, heading_style='ATX', code_language='')


def detect_api_base_path():
    global api_base_path
    for path in ['/rest/api', '/wiki/rest/api']:
        try:
            url = f'{CONFLUENCE_URL}{path}/space?limit=1'
            kwargs = {'headers': get_headers(), 'timeout': 10}
            auth = get_auth()
            if auth:
                kwargs['auth'] = auth
            resp = http_req.get(url, **kwargs)
            if resp.ok:
                api_base_path = path
                print(f'Detected API base path: {path}')
                return
        except Exception:
            pass
    print('Could not auto-detect API base path, using /rest/api')


# ===== Confluence Tool Functions (matching mcp-atlassian format) =====

def _build_page_result(data, include_content=True, convert_to_markdown=True):
    """Build a page result object matching mcp-atlassian output format."""
    space = data.get('space', {})
    version = data.get('version', {})
    by = version.get('by', {})
    links = data.get('_links', {})
    web_url = f"{CONFLUENCE_URL}{links.get('webui', '')}" if links.get('webui') else ''
    labels = [l.get('name') for l in data.get('metadata', {}).get('labels', {}).get('results', [])]

    result = {
        'id': data.get('id'),
        'title': data.get('title'),
        'type': data.get('type', 'page'),
        'url': web_url,
        'space': {'key': space.get('key', ''), 'name': space.get('name', '')},
        'version': version.get('number'),
        'updated': version.get('when', ''),
        'updated_by': by.get('displayName') or by.get('username', ''),
        'labels': labels,
    }

    if include_content:
        body_html = data.get('body', {}).get('storage', {}).get('value', '')
        if convert_to_markdown:
            result['content'] = {'value': html_to_markdown(body_html)[:8000], 'format': 'view'}
        else:
            result['content'] = {'value': body_html[:8000], 'format': 'storage'}

    return result


def _forced_search_space_key(space_key):
    if space_key is None:
        return None
    s = str(space_key).strip()
    return s if s else None


def _parse_force_space_keys(space_key):
    """UI에서 쉼표로 이어 보낸 다중 스페이스 키. None이면 스코프 제한 없음."""
    fk = _forced_search_space_key(space_key)
    if not fk:
        return None
    parts = [x.strip() for x in fk.split(',') if x.strip()]
    return frozenset(parts) if parts else None


def _cql_space_predicate_from_force(force_space_key):
    """단일/다중 스페이스 스코프용 AND (...space...) 절. 없으면 빈 문자열."""
    fk = _forced_search_space_key(force_space_key)
    if not fk:
        return ''
    parts = [x.strip() for x in fk.split(',') if x.strip()]
    if not parts:
        return ''
    if len(parts) == 1:
        return f' AND space="{parts[0]}"'
    cond = ' OR '.join(f'space="{p}"' for p in parts)
    return f' AND ({cond})'


def _scope_label_for_errors(force_space_key):
    allowed = _parse_force_space_keys(force_space_key)
    if allowed is None:
        return ''
    if len(allowed) == 1:
        return next(iter(allowed))
    return f'{len(allowed)}개 스페이스'


def _content_space_key(data):
    if not isinstance(data, dict):
        return None
    sp = data.get('space') or {}
    return sp.get('key') if isinstance(sp, dict) else None


def _strip_cql_space_predicates(cql: str) -> str:
    """Remove AND space="..." / AND space='...' so a single enforced space can be applied."""
    s = cql.strip()
    prev = None
    while prev != s:
        prev = s
        s = re.sub(r'\s+AND\s+space\s*=\s*"[^"]*"', '', s, flags=re.I)
        s = re.sub(r"\s+AND\s+space\s*=\s*'[^']*'", '', s, flags=re.I)
    s = re.sub(r'^\s*space\s*=\s*"[^"]*"\s+AND\s+', '', s, flags=re.I)
    s = re.sub(r"^\s*space\s*=\s*'[^']*'\s+AND\s+", '', s, flags=re.I)
    return s.strip()


def tool_search(query, limit=10, spaces_filter=None, force_space_key=None):
    if any(c in query for c in ['=', '~', '>', '<']):
        cql = query
    else:
        cql = f'siteSearch ~ "{query}"'
    if 'type' not in cql.lower():
        cql += ' AND type=page'

    fk = (force_space_key or '').strip()
    if fk:
        cql = _strip_cql_space_predicates(cql) + _cql_space_predicate_from_force(force_space_key)
    elif spaces_filter and not re.search(r'\bspace\s*=', cql, re.I):
        space_keys = [s.strip() for s in spaces_filter.split(',') if s.strip()]
        if len(space_keys) == 1:
            cql += f' AND space="{space_keys[0]}"'
        elif space_keys:
            space_cond = ' OR '.join(f'space="{k}"' for k in space_keys)
            cql += f' AND ({space_cond})'

    def parse(data):
        results = []
        for r in data.get('results', []):
            content = r.get('content') or {}
            space_info = r.get('resultGlobalContainer') or {}
            links = content.get('_links', {})
            web_url = f"{CONFLUENCE_URL}{links.get('webui', '')}" if links.get('webui') else ''
            title = content.get('title') or r.get('title', '')
            title = title.replace('@@@hl@@@', '').replace('@@@endhl@@@', '')
            excerpt = (r.get('excerpt', '') or '').replace('@@@hl@@@', '').replace('@@@endhl@@@', '')

            item = {
                'id': content.get('id'),
                'title': title,
                'type': content.get('type', 'page'),
                'created': '',
                'updated': r.get('friendlyLastModified') or r.get('lastModified', ''),
                'url': web_url,
                'space': {'key': space_info.get('displayUrl', '').split('/display/')[-1].split('/')[0] if space_info.get('displayUrl') else '',
                          'name': space_info.get('title', '')},
                'content': {'value': excerpt[:500], 'format': 'view'} if excerpt else {'value': '', 'format': 'view'},
            }
            results.append(item)
        return results

    try:
        endpoint = f'/search?cql={http_req.utils.quote(cql)}&limit={limit}'
        data = confluence_api('GET', endpoint)
        return json.dumps(parse(data), ensure_ascii=False)
    except Exception:
        try:
            fk = (force_space_key or '').strip()
            simple_kw = not any(c in query for c in ['=', '~', '>', '<'])
            if fk:
                scope = _cql_space_predicate_from_force(force_space_key)
                if simple_kw:
                    cql2 = f'(text ~ "{query}" OR title ~ "{query}") AND type=page' + scope
                else:
                    cql2 = _strip_cql_space_predicates(cql) + scope
            else:
                cql2 = f'text ~ "{query}" OR title ~ "{query}"'
                if 'type' not in cql2.lower():
                    cql2 += ' AND type=page'
                if spaces_filter:
                    space_keys = [s.strip() for s in spaces_filter.split(',') if s.strip()]
                    if len(space_keys) == 1:
                        cql2 = f'({cql2}) AND space="{space_keys[0]}"'
                    elif space_keys:
                        space_cond = ' OR '.join(f'space="{k}"' for k in space_keys)
                        cql2 = f'({cql2}) AND ({space_cond})'
            endpoint = f'/search?cql={http_req.utils.quote(cql2)}&limit={limit}'
            data = confluence_api('GET', endpoint)
            return json.dumps(parse(data), ensure_ascii=False)
        except Exception as e:
            return json.dumps({'error': str(e)})


def tool_get_page(page_id=None, title=None, space_key=None, include_metadata=True, convert_to_markdown=True, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if page_id:
            data = confluence_api('GET', f'/content/{page_id}?expand=body.storage,version,space,ancestors,metadata.labels')
        elif title and space_key:
            sk_arg = str(space_key).strip()
            if allowed is not None and sk_arg not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps(
                    {
                        'error': (
                            f'이 채팅은 스페이스 범위 "{lab}" 안에서만 검색할 수 있습니다. '
                            f'요청한 space_key "{space_key}"는 그 범위에 없습니다.'
                        )
                    },
                    ensure_ascii=False,
                )
            cql = f'title="{title}" AND space="{space_key}"'
            search_data = confluence_api('GET', f'/search?cql={http_req.utils.quote(cql)}&limit=1')
            results = search_data.get('results', [])
            if not results:
                return json.dumps({'error': f'Page not found: "{title}" in space {space_key}'})
            content = results[0].get('content', {})
            pid = content.get('id')
            data = confluence_api('GET', f'/content/{pid}?expand=body.storage,version,space,ancestors,metadata.labels')
        else:
            return json.dumps({'error': 'Provide page_id, or both title and space_key'})

        if allowed is not None:
            pk = _content_space_key(data)
            if pk not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps(
                    {
                        'error': (
                            f'페이지는 스페이스 "{pk}"에 있으며, 현재 채팅 스코프("{lab}") 밖입니다.'
                        )
                    },
                    ensure_ascii=False,
                )

        return json.dumps(_build_page_result(data, include_content=True, convert_to_markdown=convert_to_markdown), ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_create_page(space_key, title, content, parent_id=None, content_format='markdown', force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    sk = str(space_key).strip()
    if allowed is not None and sk not in allowed:
        lab = _scope_label_for_errors(force_space_key)
        return json.dumps({'error': f'이 채팅은 스페이스 범위 "{lab}" 안에만 글을 쓸 수 있습니다.'}, ensure_ascii=False)
    try:
        if allowed is not None and parent_id:
            par = confluence_api('GET', f'/content/{parent_id}?expand=space')
            if _content_space_key(par) not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps({'error': f'부모 페이지가 스코프 스페이스 범위("{lab}")에 있지 않습니다.'}, ensure_ascii=False)
        if content_format == 'markdown':
            from markdownify import markdownify
            import markdown as md_lib
            try:
                html_content = md_lib.markdown(content, extensions=['tables', 'fenced_code'])
            except Exception:
                html_content = f'<p>{content}</p>'
        else:
            html_content = content

        payload = {
            'type': 'page', 'title': title, 'space': {'key': space_key},
            'body': {'storage': {'value': html_content, 'representation': 'storage'}},
        }
        if parent_id:
            payload['ancestors'] = [{'id': str(parent_id)}]
        data = confluence_api('POST', '/content', payload)
        return json.dumps(_build_page_result(data, include_content=False), ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_update_page(page_id, title, content, is_minor_edit=False, version_comment=None, content_format='markdown', force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        current = confluence_api('GET', f'/content/{page_id}?expand=version,space')
        if allowed is not None and _content_space_key(current) not in allowed:
            lab = _scope_label_for_errors(force_space_key)
            return json.dumps({'error': f'페이지가 스코프 스페이스("{lab}")에 있지 않아 수정할 수 없습니다.'}, ensure_ascii=False)

        if content_format == 'markdown':
            try:
                import markdown as md_lib
                html_content = md_lib.markdown(content, extensions=['tables', 'fenced_code'])
            except Exception:
                html_content = f'<p>{content}</p>'
        else:
            html_content = content

        version_info = {'number': current['version']['number'] + 1, 'minorEdit': is_minor_edit}
        if version_comment:
            version_info['message'] = version_comment

        payload = {
            'type': 'page', 'title': title,
            'body': {'storage': {'value': html_content, 'representation': 'storage'}},
            'version': version_info,
        }
        data = confluence_api('PUT', f'/content/{page_id}', payload)
        return json.dumps(_build_page_result(data, include_content=False), ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_delete_page(page_id, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if allowed is not None:
            cur = confluence_api('GET', f'/content/{page_id}?expand=space')
            if _content_space_key(cur) not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps({'error': f'페이지가 스코프 스페이스("{lab}")에 있지 않아 삭제할 수 없습니다.'}, ensure_ascii=False)
        confluence_api('DELETE', f'/content/{page_id}')
        return json.dumps({'success': True, 'message': f'Page {page_id} deleted.'})
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_get_spaces(limit=50, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if allowed is not None:
            out = []
            for k in sorted(allowed):
                try:
                    data = confluence_api('GET', f'/space/{http_req.utils.quote(k, safe="~.")}')
                    out.append({'key': data.get('key'), 'name': data.get('name')})
                except Exception:
                    out.append({'key': k, 'name': k})
            return json.dumps(out, ensure_ascii=False)
        data = confluence_api('GET', f'/space?limit={limit}&type=global')
        spaces = [{'key': s.get('key'), 'name': s.get('name')} for s in data.get('results', [])]
        return json.dumps(spaces, ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_get_comments(page_id, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if allowed is not None:
            pg = confluence_api('GET', f'/content/{page_id}?expand=space')
            if _content_space_key(pg) not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps({'error': f'댓글 대상 페이지가 스코프 스페이스("{lab}")에 있지 않습니다.'}, ensure_ascii=False)
        data = confluence_api('GET', f'/content/{page_id}/child/comment?expand=body.storage,version')
        comments = []
        for c in data.get('results', []):
            v = c.get('version', {})
            by = v.get('by', {})
            comments.append({
                'id': c.get('id'),
                'author': by.get('displayName') or by.get('username', 'Unknown'),
                'body': html_to_markdown(c.get('body', {}).get('storage', {}).get('value', '')),
                'created': v.get('when', ''),
            })
        return json.dumps(comments, ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_add_comment(page_id, body, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if allowed is not None:
            pg = confluence_api('GET', f'/content/{page_id}?expand=space')
            if _content_space_key(pg) not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps({'error': f'댓글을 달 페이지가 스코프 스페이스("{lab}")에 있지 않습니다.'}, ensure_ascii=False)
        data = confluence_api('POST', '/content', {
            'type': 'comment', 'container': {'id': str(page_id), 'type': 'page'},
            'body': {'storage': {'value': f'<p>{body}</p>', 'representation': 'storage'}},
        })
        return json.dumps({'id': data.get('id'), 'success': True})
    except Exception as e:
        return json.dumps({'error': str(e)})


def tool_get_children(parent_id, limit=25, force_space_key=None):
    allowed = _parse_force_space_keys(force_space_key)
    try:
        if allowed is not None:
            parent = confluence_api('GET', f'/content/{parent_id}?expand=space')
            if _content_space_key(parent) not in allowed:
                lab = _scope_label_for_errors(force_space_key)
                return json.dumps({'error': f'부모 페이지가 스코프 스페이스("{lab}")에 있지 않아 하위 목록을 볼 수 없습니다.'}, ensure_ascii=False)
        data = confluence_api('GET', f'/content/{parent_id}/child/page?expand=version&limit={limit}')
        children = []
        for c in data.get('results', []):
            v = c.get('version', {})
            links = c.get('_links', {})
            web_url = f"{CONFLUENCE_URL}{links.get('webui', '')}" if links.get('webui') else ''
            children.append({
                'id': c.get('id'), 'title': c.get('title'),
                'type': 'page', 'url': web_url,
                'version': v.get('number'), 'updated': v.get('when', ''),
            })
        return json.dumps(children, ensure_ascii=False)
    except Exception as e:
        return json.dumps({'error': str(e)})


TOOL_MAP = {
    'confluence_search': tool_search, 'confluence_get_page': tool_get_page,
    'confluence_create_page': tool_create_page, 'confluence_update_page': tool_update_page,
    'confluence_delete_page': tool_delete_page,
    'confluence_get_spaces': tool_get_spaces,
    'confluence_get_comments': tool_get_comments, 'confluence_add_comment': tool_add_comment,
    'confluence_get_page_children': tool_get_children,
}

TOOLS = [
    {
        "name": "confluence_search",
        "description": "Search Confluence content using simple terms or CQL.\nWhen the chat has a target space selected in the UI, the server forces this search to that space only.\nReturns: JSON list of simplified Confluence page objects.",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Search query - can be either a simple text (e.g. 'project documentation') or a CQL query string. Simple queries use 'siteSearch' by default. Examples of CQL:\n- Basic search: 'type=page AND space=DEV'\n- Search by title: 'title~\"Meeting Notes\"'\n- Recent content: 'created >= \"2023-01-01\"'\n- Content with specific label: 'label=documentation'\n- Recently modified content: 'lastModified > startOfMonth(\"-1M\")'\n- Content modified this year: 'creator = currentUser() AND lastModified > startOfYear()'\n- Title wildcards: 'title ~ \"Minutes*\" AND (space = \"HR\" OR space = \"Marketing\")'",
                },
                "limit": {"type": "integer", "description": "Maximum number of results (1-50)", "default": 10, "minimum": 1, "maximum": 50},
                "spaces_filter": {"type": "string", "description": "(Optional) Comma-separated list of space keys to filter results by."},
            },
            "required": ["query"],
        },
    },
    {
        "name": "confluence_get_page",
        "description": "Get content of a specific Confluence page by its ID, or by its title and space key.\nWhen a target space is selected in the UI, the server only allows pages in that space.\nReturns: JSON string representing the page content and metadata.",
        "parameters": {
            "type": "object",
            "properties": {
                "page_id": {"type": "string", "description": "Confluence page ID (numeric ID from URL). Provide this OR both 'title' and 'space_key'."},
                "title": {"type": "string", "description": "The exact title of the page. Use with 'space_key' if 'page_id' is not known."},
                "space_key": {"type": "string", "description": "The key of the space. Required if using 'title'."},
                "include_metadata": {"type": "boolean", "description": "Whether to include page metadata.", "default": True},
                "convert_to_markdown": {"type": "boolean", "description": "Convert content to markdown (true) or keep raw HTML (false).", "default": True},
            },
        },
    },
    {
        "name": "confluence_create_page",
        "description": "Create a new Confluence page. Resolve the parent first (UI picker, pasted URL with /pages/<id>/, confluence_search, or confluence_get_page). Omit parent_id only when the user explicitly wants the page at the space root.\nReturns: JSON with id, title, url, space, etc.",
        "parameters": {
            "type": "object",
            "properties": {
                "space_key": {"type": "string", "description": "The key of the space to create the page in (e.g. 'DEV', 'TEAM')"},
                "title": {"type": "string", "description": "The title of the page"},
                "content": {"type": "string", "description": "The content of the page. Format depends on content_format parameter."},
                "parent_id": {"type": "string", "description": "Optional. Numeric Confluence page ID for the parent. Use UI selection, URL extraction, confluence_search, or confluence_get_page to obtain it. Omit only for explicit space-root pages."},
                "content_format": {"type": "string", "description": "(Optional) 'markdown' (default), 'wiki', or 'storage'", "default": "markdown"},
            },
            "required": ["space_key", "title", "content"],
        },
    },
    {
        "name": "confluence_update_page",
        "description": "Update an existing Confluence page. Resolve page_id from UI selection, pasted URL, confluence_search, or confluence_get_page.\nReturns: JSON string representing the updated page object.",
        "parameters": {
            "type": "object",
            "properties": {
                "page_id": {"type": "string", "description": "Numeric page ID. Prefer the UI-selected target; otherwise URL extraction or search/get_page."},
                "title": {"type": "string", "description": "The new title of the page"},
                "content": {"type": "string", "description": "The new content of the page"},
                "is_minor_edit": {"type": "boolean", "description": "Whether this is a minor edit", "default": False},
                "version_comment": {"type": "string", "description": "Optional comment for this version"},
                "content_format": {"type": "string", "description": "(Optional) 'markdown' (default), 'wiki', or 'storage'", "default": "markdown"},
            },
            "required": ["page_id", "title", "content"],
        },
    },
    {
        "name": "confluence_delete_page",
        "description": "Delete an existing Confluence page. Resolve page_id from UI selection, pasted URL, confluence_search, or confluence_get_page.\nReturns: JSON string indicating success or failure.",
        "parameters": {
            "type": "object",
            "properties": {"page_id": {"type": "string", "description": "Numeric page ID. Prefer the UI-selected target; otherwise URL extraction or search/get_page."}},
            "required": ["page_id"],
        },
    },
    {
        "name": "confluence_get_spaces",
        "description": "List Confluence spaces.\nWhen a target space is selected in the UI, returns only that space.\nReturns: JSON list of space objects with key and name.",
        "parameters": {
            "type": "object",
            "properties": {"limit": {"type": "integer", "description": "Max spaces to return", "default": 50}},
        },
    },
    {
        "name": "confluence_get_comments",
        "description": "Get comments for a specific Confluence page.\nReturns: JSON list of comment objects.",
        "parameters": {
            "type": "object",
            "properties": {"page_id": {"type": "string", "description": "Confluence page ID"}},
            "required": ["page_id"],
        },
    },
    {
        "name": "confluence_add_comment",
        "description": "Add a comment to a Confluence page. Resolve page_id from a pasted URL (.../pages/<id>/...), confluence_search, or confluence_get_page when the user has not picked a page in the UI.\nReturns: JSON with id and success.",
        "parameters": {
            "type": "object",
            "properties": {
                "page_id": {"type": "string", "description": "Numeric Confluence page ID. Use UI-selected target, URL extraction, confluence_search, or confluence_get_page."},
                "body": {"type": "string", "description": "The comment content (plain text; server wraps in storage format)."},
            },
            "required": ["page_id", "body"],
        },
    },
    {
        "name": "confluence_get_page_children",
        "description": "Get child pages of a specific Confluence page.\nReturns: JSON list of child page objects.",
        "parameters": {
            "type": "object",
            "properties": {
                "parent_id": {"type": "string", "description": "The ID of the parent page"},
                "limit": {"type": "integer", "description": "Maximum number of child items (1-50)", "default": 25},
            },
            "required": ["parent_id"],
        },
    },
]


def filter_tools(allowed=None):
    if not allowed:
        return TOOLS
    return [t for t in TOOLS if t["name"] in allowed]


def tools_to_openai(allowed=None):
    return [{"type": "function", "function": {"name": t["name"], "description": t["description"], "parameters": t["parameters"]}} for t in filter_tools(allowed)]



BASE_SYSTEM_PROMPT = f"""You are a Confluence assistant for {CONFLUENCE_URL}. Respond in the user's language.
When presenting tool results, always include the url field as a clickable markdown link. Do not omit or rephrase fields from tool results."""


WRITE_TOOL_LABELS = {
    'confluence_create_page': '페이지 생성',
    'confluence_update_page': '페이지 수정',
    'confluence_delete_page': '페이지 삭제',
    'confluence_add_comment': '댓글 작성',
}

READ_ONLY_TOOL_NAMES = [
    'confluence_search',
    'confluence_get_page',
    'confluence_get_spaces',
    'confluence_get_page_children',
    'confluence_get_comments',
]
SAFE_WRITE_TOOL_NAMES = ['confluence_create_page', 'confluence_add_comment']


def normalize_enabled_tools(enabled_tools):
    """If the client omits enabled_tools, never expose destructive tools by default."""
    if enabled_tools is None:
        return READ_ONLY_TOOL_NAMES + SAFE_WRITE_TOOL_NAMES
    return enabled_tools


def extract_confluence_page_id(url_or_id):
    """Accept a numeric id or a Confluence URL; return page id string or None."""
    if not url_or_id or not isinstance(url_or_id, str):
        return None
    s = url_or_id.strip()
    if s.isdigit():
        return s
    m = re.search(r'/pages/(\d+)', s)
    if m:
        return m.group(1)
    m = re.search(r'pageId=(\d+)', s, re.I)
    if m:
        return m.group(1)
    m = re.search(r'/content/(\d+)', s)
    if m:
        return m.group(1)
    return None


def _page_creation_guidance(allowed_tools, parent_page_context):
    """Extra system text for create-page workflow, confirmation, and optional UI parent."""
    if allowed_tools is not None and 'confluence_create_page' not in allowed_tools:
        return ''
    lines = [
        'Page creation (confluence_create_page): Users rarely know numeric page IDs. If they paste a Confluence URL, '
        'extract the page ID from a path like .../pages/<digits>/... (or pageId= in query). Otherwise use '
        'confluence_search or confluence_get_page (title + space_key) to resolve the parent before creating. '
        'Only omit parent_id when the user explicitly wants the new page at the space root.',
        'After a successful confluence_create_page tool result: state the new page title, space key, whether it sits '
        'under a parent (name the parent) or at root, include the result url as a markdown link, and ask the user '
        'to confirm the location is correct. If the location is wrong, suggest fixing it in the Confluence UI or creating again under the correct parent.',
    ]
    if isinstance(parent_page_context, dict):
        pid = str(parent_page_context.get('id') or '').strip()
        ptitle = str(parent_page_context.get('title') or '').strip()
        psk = str(parent_page_context.get('space_key') or '').strip()
        if pid:
            lines.append(
                f'The user selected a default parent page in the UI for this turn: parent_id="{pid}", '
                f'title="{ptitle}", space_key="{psk}". Use this parent_id in confluence_create_page unless the user '
                f'clearly specifies a different parent or explicitly requests a root-level page.'
            )
    return '\n' + ' '.join(lines)


def _comment_target_guidance(allowed_tools, comment_target_page_context):
    """UI-selected page for confluence_add_comment and workflow hints."""
    if allowed_tools is not None and 'confluence_add_comment' not in allowed_tools:
        return ''
    lines = [
        'Comments (confluence_add_comment): Resolve page_id from a pasted Confluence URL, confluence_search, or '
        'confluence_get_page before posting when the user has not selected a page in the UI.',
        'After a successful confluence_add_comment, briefly confirm which page (title and space) received the comment.',
    ]
    if isinstance(comment_target_page_context, dict):
        pid = str(comment_target_page_context.get('id') or '').strip()
        ptitle = str(comment_target_page_context.get('title') or '').strip()
        psk = str(comment_target_page_context.get('space_key') or '').strip()
        if pid:
            lines.append(
                f'The user selected a target page in the UI for comments: page_id="{pid}", title="{ptitle}", '
                f'space_key="{psk}". Use this page_id in confluence_add_comment unless the user clearly specifies a different page.'
            )
    return '\n' + ' '.join(lines)


def _update_target_guidance(allowed_tools, update_target_page_context):
    if allowed_tools is not None and 'confluence_update_page' not in allowed_tools:
        return ''
    lines = [
        'Page updates (confluence_update_page): Resolve page_id via UI selection, pasted URL, confluence_search, or confluence_get_page. '
        'After success, confirm which page was updated (title, space).',
    ]
    if isinstance(update_target_page_context, dict):
        pid = str(update_target_page_context.get('id') or '').strip()
        ptitle = str(update_target_page_context.get('title') or '').strip()
        psk = str(update_target_page_context.get('space_key') or '').strip()
        if pid:
            lines.append(
                f'The user selected the page to edit in the UI: page_id="{pid}", title="{ptitle}", space_key="{psk}". '
                f'Use this page_id in confluence_update_page unless the user clearly names a different page.'
            )
    return '\n' + ' '.join(lines)


def _delete_target_guidance(allowed_tools, delete_target_page_context):
    if allowed_tools is not None and 'confluence_delete_page' not in allowed_tools:
        return ''
    lines = [
        'Page deletion (confluence_delete_page): Resolve page_id via UI selection, pasted URL, confluence_search, or confluence_get_page. '
        'After success, confirm which page was deleted.',
    ]
    if isinstance(delete_target_page_context, dict):
        pid = str(delete_target_page_context.get('id') or '').strip()
        ptitle = str(delete_target_page_context.get('title') or '').strip()
        psk = str(delete_target_page_context.get('space_key') or '').strip()
        if pid:
            lines.append(
                f'The user selected the page to delete in the UI: page_id="{pid}", title="{ptitle}", space_key="{psk}". '
                f'Use this page_id in confluence_delete_page unless the user clearly names a different page.'
            )
    return '\n' + ' '.join(lines)


def get_system_prompt(
    space_key=None,
    allowed_tools=None,
    parent_page_context=None,
    comment_target_page_context=None,
    update_target_page_context=None,
    delete_target_page_context=None,
):
    prompt = BASE_SYSTEM_PROMPT
    if space_key:
        sk = str(space_key).strip()
        nkeys = len([p for p in sk.split(',') if p.strip()])
        if nkeys <= 1:
            prompt += (
                f'\nTarget space: "{sk}". The server restricts all Confluence tools (search, get page, children, comments, spaces list, and writes) '
                f'to this space only; you cannot load or modify content from other spaces.'
            )
        else:
            prompt += (
                f'\nTarget scope: {nkeys} Confluence spaces (comma-separated keys from the UI). The server restricts all Confluence tools '
                f'to those spaces only; you cannot load or modify content from spaces outside that set.'
            )
    if allowed_tools is not None:
        disabled = [label for name, label in WRITE_TOOL_LABELS.items() if name not in allowed_tools]
        if disabled:
            names = ', '.join(disabled)
            prompt += f'\nDisabled: {names}. Do not attempt these. Tell the user to pick that operation in the chat task bar (only one write task can be active per message).'
    prompt += (
        '\nWhen the user asks to create, comment on, update, or delete Confluence content, call the appropriate tool as usual. '
        'The server shows the user a preview of the action and runs it only after they confirm in the UI; do not state that the change is already done until they have confirmed.'
    )
    prompt += (
        '\nThe user may attach files (documents, spreadsheets, slides, PDFs, images). '
        'Extracted text and any images in the message are part of their request—answer using both their words and the attachment content.'
    )
    prompt += (
        '\nUploaded-file priority: If the user message contains "### 첨부 파일:" blocks or long pasted document text from an upload, '
        'treat that as the primary source for phrases like "this paper", "this file", "the attachment", "첨부", "위 논문", or "이 문서". '
        'Summarize or answer from that content first. Do not call confluence_search (or other tools) to hunt for a Confluence page '
        'unless the user clearly asks about Confluence spaces/pages by name or URL.'
    )
    prompt += (
        '\nConfluence writes from uploads: When the user wants to publish, post, summarize on Confluence, create a page, add a comment, '
        'or update a page using data from files they attached (or text you already extracted in this conversation), synthesize from that '
        'attachment-derived analysis. Call confluence_create_page, confluence_add_comment, or confluence_update_page with accurate markdown '
        'or storage-format body as required. Do not rely on confluence_search to rediscover the same source unless they name a specific page. '
        'If a write tool is disabled in the UI, tell them to enable that task (e.g. page creation) in the chat toolbar.'
    )
    prompt += _page_creation_guidance(allowed_tools, parent_page_context)
    prompt += _comment_target_guidance(allowed_tools, comment_target_page_context)
    prompt += _update_target_guidance(allowed_tools, update_target_page_context)
    prompt += _delete_target_guidance(allowed_tools, delete_target_page_context)
    return prompt


_SCOPED_TOOL_NAMES = (
    'confluence_search',
    'confluence_get_page',
    'confluence_get_page_children',
    'confluence_get_comments',
    'confluence_get_spaces',
    'confluence_create_page',
    'confluence_update_page',
    'confluence_delete_page',
    'confluence_add_comment',
)

# 쓰기 도구는 서버가 한 번 더 미리보기를 보여주고 사용자 확인 후에만 실행합니다.
TOOLS_REQUIRING_CONFIRMATION = frozenset({
    'confluence_create_page',
    'confluence_add_comment',
    'confluence_update_page',
    'confluence_delete_page',
})


def _truncate_preview_text(s, max_len=2800):
    s = s if s is not None else ''
    s = str(s)
    if len(s) <= max_len:
        return s
    return s[:max_len] + '\n\n…_(이하 생략)_'


def _preview_non_write_tool(name, args):
    args = args or {}
    if name == 'confluence_search':
        return f"- **검색** — 쿼리: `{args.get('query', '')}`"
    if name == 'confluence_get_page':
        parts = []
        if args.get('page_id'):
            parts.append(f"page_id `{args['page_id']}`")
        if args.get('title'):
            parts.append(f"제목 `{args['title']}`")
        if args.get('space_key'):
            parts.append(f"스페이스 `{args['space_key']}`")
        return f"- **페이지 조회** — {', '.join(parts) or '인자 없음'}"
    if name == 'confluence_get_page_children':
        return f"- **하위 페이지** — parent_id `{args.get('parent_id', '')}`"
    if name == 'confluence_get_comments':
        return f"- **댓글 목록** — page_id `{args.get('page_id', '')}`"
    if name == 'confluence_get_spaces':
        return '- **스페이스 목록**'
    brief = json.dumps(args, ensure_ascii=False)
    if len(brief) > 400:
        brief = brief[:400] + '…'
    return f"- **`{name}`** — {brief}"


def _format_delete_preview_from_page_data(data, page_id_str):
    """Confluence content JSON으로 삭제 확인용 마크다운 블록 생성."""
    space = data.get('space') or {}
    version = data.get('version') or {}
    by = version.get('by') or {}
    links = data.get('_links') or {}
    web_url = f"{CONFLUENCE_URL}{links.get('webui', '')}" if links.get('webui') else ''
    title = data.get('title') or '_(제목 없음)_'
    vc = version.get('number')
    when = version.get('when') or ''
    editor = (
        by.get('displayName')
        or by.get('username')
        or by.get('publicName')
        or '—'
    )

    ancestors = data.get('ancestors') or []
    path_parts = [a.get('title', '') for a in ancestors if isinstance(a, dict) and a.get('title')]
    path_line = ' > '.join(path_parts) if path_parts else None

    labels = []
    meta = data.get('metadata') or {}
    lab = meta.get('labels') or {}
    for l in lab.get('results') or []:
        if isinstance(l, dict) and l.get('name'):
            labels.append(l['name'])

    hist = data.get('history') or {}
    created = hist.get('createdDate') or ''
    created_by = hist.get('createdBy') or {}
    if not isinstance(created_by, dict):
        created_by = {}
    creator = (
        created_by.get('displayName')
        or created_by.get('username')
        or created_by.get('publicName')
        or ''
    )

    lines = [
        f'- **페이지 ID**: `{page_id_str}`',
        f'- **제목**: {title}',
        f'- **스페이스**: `{space.get("key", "")}` ({space.get("name", "") or "—"})',
    ]
    if web_url:
        lines.append(f'- **URL**: [{web_url}]({web_url})')
    lines.append(f'- **현재 버전**: {vc if vc is not None else "—"}')
    lines.append(f'- **최종 수정 시각**: {when or "—"}')
    lines.append(f'- **최종 수정자**: {editor}')
    if created or creator:
        created_line = f'- **최초 작성 시각**: {created or "—"}'
        if creator:
            created_line += f' — **작성자**: {creator}'
        lines.append(created_line)
    if path_line:
        lines.append(f'- **상위 경로**: {path_line}')
    if labels:
        lines.append(f'- **라벨**: {", ".join(labels)}')
    return '\n'.join(lines)


def _delete_preview_detail_markdown(page_id, force_space_key=None):
    """삭제 확인 전 실제 문서 메타를 조회해 표시 (실패 시 오류 안내 + ID만)."""
    if not page_id or not str(page_id).strip():
        return None
    pid = str(page_id).strip()
    allowed = _parse_force_space_keys(force_space_key)
    scope_lab = _scope_label_for_errors(force_space_key)
    data = None
    try:
        data = confluence_api(
            'GET',
            f'/content/{pid}?expand=version,space,ancestors,metadata.labels,history',
        )
    except Exception:
        try:
            data = confluence_api(
                'GET',
                f'/content/{pid}?expand=version,space,ancestors,metadata.labels',
            )
        except Exception as e:
            return (
                f'_페이지 정보를 불러오지 못했습니다: {e}_\n\n'
                f'- **페이지 ID**: `{pid}`'
            )

    if allowed is not None and _content_space_key(data) not in allowed:
        pk = _content_space_key(data) or '—'
        warn = (
            f'_⚠️ 이 페이지는 스페이스 `{pk}`에 있으며, 현재 채팅 스코프(`{scope_lab}`)와 다를 수 있습니다. '
            f'실행 시 서버에서 거절될 수 있습니다._\n\n'
        )
        return warn + _format_delete_preview_from_page_data(data, pid)

    return _format_delete_preview_from_page_data(data, pid)


def format_write_action_preview(name, args, force_space_key=None):
    args = args or {}
    if name == 'confluence_create_page':
        parent = args.get('parent_id')
        parent_line = f"- **부모 페이지 ID**: `{parent}`" if parent else '- **부모**: (루트 또는 스페이스 최상위)'
        body = _truncate_preview_text(args.get('content') or '')
        fmt = args.get('content_format') or 'markdown'
        return (
            '### 페이지 생성 예정\n'
            f"- **스페이스 키**: `{args.get('space_key', '')}`\n"
            f"- **제목**: {args.get('title', '') or '_(없음)_'}\n"
            f"{parent_line}\n"
            f"- **본문 형식**: `{fmt}`\n\n"
            '**본문 미리보기**\n\n'
            f'{body}'
        )
    if name == 'confluence_add_comment':
        body = _truncate_preview_text(args.get('body') or '')
        return (
            '### 댓글 작성 예정\n'
            f"- **대상 페이지 ID**: `{args.get('page_id', '')}`\n\n"
            '**댓글 내용**\n\n'
            f'{body}'
        )
    if name == 'confluence_update_page':
        body = _truncate_preview_text(args.get('content') or '')
        extra = []
        if args.get('is_minor_edit'):
            extra.append('사소한 수정')
        if args.get('version_comment'):
            extra.append(f"버전 코멘트: {args['version_comment']}")
        extra_s = (' — ' + ', '.join(extra)) if extra else ''
        return (
            '### 페이지 수정 예정\n'
            f"- **페이지 ID**: `{args.get('page_id', '')}`\n"
            f"- **새 제목**: {args.get('title', '') or '_(변경 없음/기존 유지)_'}{extra_s}\n\n"
            '**바꿀 본문 미리보기**\n\n'
            f'{body}'
        )
    if name == 'confluence_delete_page':
        detail = _delete_preview_detail_markdown(args.get('page_id'), force_space_key)
        body = detail if detail else f"- **페이지 ID**: `{args.get('page_id', '') or '(없음)'}`"
        return (
            '### 페이지 삭제 예정\n\n'
            f'{body}\n\n'
            '⚠️ 삭제는 되돌리기 어렵습니다. 계속하려면 아래에서 실행을 확인하세요.'
        )
    return f'### `{name}`\n\n```json\n{json.dumps(args, ensure_ascii=False, indent=2)[:2000]}\n```'


def format_pending_confirmation_markdown(calls, force_space_key=None):
    """calls: list of {"name", "arguments"}"""
    blocks = []
    for i, c in enumerate(calls, start=1):
        name = c.get('name') or ''
        args = c.get('arguments') if isinstance(c.get('arguments'), dict) else {}
        if name in TOOLS_REQUIRING_CONFIRMATION:
            blocks.append(format_write_action_preview(name, args, force_space_key))
        else:
            blocks.append(f'### 함께 실행될 작업 ({i})\n\n{_preview_non_write_tool(name, args)}')
    return '\n\n---\n\n'.join(blocks)


def execute_tool(name, input_data, forced_scope_space_key=None):
    input_data = dict(input_data or {})
    input_data.pop('force_space_key', None)
    fsk = _forced_search_space_key(forced_scope_space_key)
    if fsk and name in _SCOPED_TOOL_NAMES:
        input_data['force_space_key'] = fsk
    func = TOOL_MAP.get(name)
    if not func:
        return json.dumps({'error': f'Unknown tool: {name}'})
    try:
        result = func(**input_data)
        print(f"\n===== TOOL CALL: {name} =====", flush=True)
        print(f"  Args: {json.dumps(input_data, ensure_ascii=False)}", flush=True)
        print(f"  Result (first 2000 chars): {result[:2000]}", flush=True)
        print(f"===== END TOOL CALL =====\n", flush=True)
        return result
    except TypeError as e:
        return json.dumps({'error': f'Invalid parameters for {name}: {str(e)}'})


# ===== Chat attachments (docx, pptx, pdf, csv, images, …) =====

def _truncate_extracted_text(text, limit):
    text = text or ''
    if len(text) <= limit:
        return text
    return text[:limit] + f'\n\n…_(추출 본문이 {limit:,}자에서 잘렸습니다.)_'


def _decode_bytes_flex(raw):
    for enc in ('utf-8', 'utf-8-sig', 'cp949', 'euc-kr', 'latin-1'):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode('utf-8', errors='replace')


def _guess_mime(filename, declared):
    ext = (os.path.splitext(filename or '')[1] or '').lower()
    if declared and declared not in ('application/octet-stream', ''):
        return declared.split(';')[0].strip().lower()
    ext_map = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.csv': 'text/csv',
        '.tsv': 'text/tab-separated-values',
        '.txt': 'text/plain',
        '.md': 'text/markdown',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.png': 'image/png',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.gif': 'image/gif',
        '.webp': 'image/webp',
        '.svg': 'image/svg+xml',
    }
    return ext_map.get(ext, declared or 'application/octet-stream')


def _extract_pdf(raw):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    parts = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text() or ''
        if t.strip():
            parts.append(f'--- PDF 페이지 {i + 1} ---\n{t.strip()}')
    return '\n\n'.join(parts) if parts else ''


def _docx_collect_from_block(block):
    """paragraphs + 표 셀 안의 단락 (1단계 표)."""
    lines = []
    for p in block.paragraphs:
        t = (p.text or '').strip()
        if t:
            lines.append(t)
    for table in block.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    t = (p.text or '').strip()
                    if t:
                        lines.append(t)
    return lines


def _extract_docx_ooxml_paragraphs(xml_bytes):
    """WordprocessingML에서 w:p 단위로 w:t만 수집 (EndNote/필드 XML 노이즈 제거)."""
    from xml.etree import ElementTree as ET

    W_NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    TAG_P = f'{{{W_NS}}}p'
    TAG_T = f'{{{W_NS}}}t'
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return ''
    paras = []
    for p in root.iter(TAG_P):
        bits = []
        for t in p.iter(TAG_T):
            if t.text:
                bits.append(t.text)
            if t.tail and t.tail.strip():
                bits.append(t.tail)
        line = ''.join(bits).strip()
        if line:
            paras.append(line)
    return '\n'.join(paras)


def _extract_docx_zip_xml(raw):
    """본문·각주·미주·주석 OOXML을 단락 텍스트로 복구."""
    import zipfile

    try:
        z = zipfile.ZipFile(io.BytesIO(raw))
    except zipfile.BadZipFile:
        return ''
    parts = []
    for name, label in (
        ('word/document.xml', 'document'),
        ('word/endnotes.xml', 'endnotes'),
        ('word/footnotes.xml', 'footnotes'),
        ('word/comments.xml', 'comments'),
    ):
        try:
            blob = z.read(name)
        except KeyError:
            continue
        txt = _extract_docx_ooxml_paragraphs(blob)
        if txt.strip():
            parts.append(f'### {label}\n{txt}')
    return '\n\n'.join(parts)


def _extract_docx(raw):
    from docx import Document

    doc = Document(io.BytesIO(raw))
    lines = _docx_collect_from_block(doc)
    try:
        for section in doc.sections:
            try:
                if section.header:
                    lines.extend(_docx_collect_from_block(section.header))
            except (AttributeError, ValueError, TypeError):
                pass
            try:
                if section.footer:
                    lines.extend(_docx_collect_from_block(section.footer))
            except (AttributeError, ValueError, TypeError):
                pass
    except (AttributeError, ValueError, TypeError):
        pass

    text = '\n'.join(lines)
    xml_fb = _extract_docx_zip_xml(raw)
    td = text.strip()
    tx = xml_fb.strip()
    if len(tx) > len(td):
        text = xml_fb
    elif len(td) < 400 and tx:
        text = (text + '\n\n' + xml_fb).strip()
    return text


def _extract_pptx(raw):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    lines = []
    for si, slide in enumerate(prs.slides, start=1):
        chunk = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                chunk.append(shape.text.strip())
        if chunk:
            lines.append(f'--- 슬라이드 {si} ---\n' + '\n'.join(chunk))
    return '\n\n'.join(lines)


def _extract_xlsx(raw):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f'### 시트: {sheet.title}')
        rows = []
        for row in sheet.iter_rows(max_row=500, max_col=50, values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(x.strip() for x in cells):
                rows.append('\t'.join(cells))
        lines.append('\n'.join(rows[:400]))
    return '\n\n'.join(lines)


def _extract_csv(raw, delimiter=None):
    text = _decode_bytes_flex(raw)
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=',\t;')
        delim = delimiter or dialect.delimiter
    except csv.Error:
        delim = delimiter or ','
    rdr = csv.reader(io.StringIO(text), delimiter=delim)
    rows = []
    for i, row in enumerate(rdr):
        if i >= 500:
            rows.append('…_(행 수 제한)_')
            break
        rows.append('\t'.join(row))
    return '\n'.join(rows)


def _extract_attachment_text(filename, mime, raw):
    """문서·표 등에서 텍스트 추출. 비전 이미지는 빈 문자열."""
    mime = (mime or '').split(';')[0].strip().lower()
    ext = (os.path.splitext(filename or '')[1] or '').lower()

    if ext == '.svg' or mime == 'image/svg+xml':
        return _decode_bytes_flex(raw)

    if mime.startswith(IMAGE_MIME_PREFIX) or ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.ico'):
        return ''

    try:
        if mime == 'application/pdf' or ext == '.pdf':
            return _extract_pdf(raw)
        if (
            mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            or ext == '.docx'
        ):
            return _extract_docx(raw)
        if (
            mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            or ext == '.pptx'
        ):
            return _extract_pptx(raw)
        if mime == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or ext == '.xlsx':
            return _extract_xlsx(raw)
        if mime in ('text/csv', 'text/tab-separated-values') or ext in ('.csv', '.tsv'):
            delim = '\t' if ext == '.tsv' or 'tab' in mime else None
            return _extract_csv(raw, delimiter=delim)
        if mime in (
            'text/plain', 'text/markdown', 'application/json', 'application/xml',
            'text/html', 'text/xml',
        ) or ext in ('.txt', '.md', '.json', '.xml', '.html', '.htm', '.log', '.yml', '.yaml', '.py', '.js', '.ts', '.css'):
            return _decode_bytes_flex(raw)
    except Exception as e:
        return f'_(이 파일에서 텍스트를 추출하지 못했습니다: {e})_'

    try:
        return _decode_bytes_flex(raw)
    except Exception:
        return f'_(지원하지 않거나 바이너리로 보이는 형식입니다: `{mime or ext}`)_'


def _build_user_message_with_attachments(user_text, attachment_items):
    """
    attachment_items: list of {filename, mime_type, raw_bytes}
    Returns dict: role, content (str | OpenAI multimodal list), display_text
    """
    user_text = (user_text or '').strip()
    names = [a.get('filename') or 'file' for a in attachment_items]

    doc_chunks = []
    image_parts = []

    for item in attachment_items:
        fn = item.get('filename') or 'attachment'
        mime = _guess_mime(fn, item.get('mime_type') or '')
        raw = item.get('raw_bytes') or b''
        if not raw:
            doc_chunks.append(f'### 첨부: {fn}\n_(빈 데이터)_')
            continue

        if mime in VISION_IMAGE_TYPES or (mime.startswith(IMAGE_MIME_PREFIX) and mime != 'image/svg+xml'):
            b64 = base64.standard_b64encode(raw).decode('ascii')
            image_parts.append({
                'type': 'image_url',
                'image_url': {'url': f'data:{mime};base64,{b64}', 'detail': 'auto'},
            })
            continue

        extracted = _extract_attachment_text(fn, mime, raw)
        extracted = _truncate_extracted_text(extracted, CHAT_ATTACHMENT_MAX_EXTRACT_CHARS)
        label = f'### 첨부 파일: {fn}\n_({mime})_\n'
        if extracted.strip():
            doc_chunks.append(label + extracted.strip())
        else:
            doc_chunks.append(label + '_(텍스트 추출 없음 — 이미지가 아니면 형식을 확인해 주세요.)_')

    display_core = user_text if user_text else '(첨부만 전송됨 — 내용을 요약·분석해 주세요.)'
    display_text = display_core
    if names:
        display_text = display_core + '\n📎 ' + ', '.join(names)

    text_body = user_text if user_text else '사용자가 파일을 첨부했습니다. 아래 추출·이미지를 바탕으로 질문에 답하거나 요약·분석해 주세요.'
    if doc_chunks:
        text_body = text_body + '\n\n---\n\n' + '\n\n---\n\n'.join(doc_chunks)

    if image_parts:
        content = [{'type': 'text', 'text': text_body}]
        content.extend(image_parts)
    else:
        content = text_body

    return {'role': 'user', 'content': content, 'display_text': display_text}


def _strip_llm_only_keys(msg):
    """응답에 포함해 클라이언트가 저장·재전송할 메시지."""
    if not isinstance(msg, dict):
        return msg
    out = {'role': msg.get('role'), 'content': msg.get('content')}
    if msg.get('display_text'):
        out['display_text'] = msg['display_text']
    if msg.get('tool_calls'):
        out['tool_calls'] = msg['tool_calls']
    if msg.get('tool_call_id'):
        out['tool_call_id'] = msg['tool_call_id']
    if msg.get('name'):
        out['name'] = msg['name']
    return out


def _message_content_for_api(msg):
    if not isinstance(msg, dict):
        return msg
    return msg.get('content')


def normalize_messages_for_chat_api(messages):
    """클라이언트 히스토리 → chat_openai용 (content만 전달)."""
    out = []
    for m in messages:
        if not isinstance(m, dict) or not m.get('role'):
            continue
        entry = {'role': m['role'], 'content': _message_content_for_api(m)}
        if m.get('tool_calls'):
            entry['tool_calls'] = m['tool_calls']
        if m.get('tool_call_id'):
            entry['tool_call_id'] = m['tool_call_id']
        out.append(entry)
    return out


def parse_chat_attachments_from_request(body):
    """
    body['attachments']: [{ filename, mime_type, data: base64 }]
    Returns (items, error_string_or_None).
    """
    raw_list = body.get('attachments')
    if not raw_list:
        return [], None
    if not isinstance(raw_list, list):
        return None, 'attachments는 배열이어야 합니다.'
    if len(raw_list) > CHAT_ATTACHMENT_MAX_FILES:
        return None, f'첨부는 최대 {CHAT_ATTACHMENT_MAX_FILES}개까지 가능합니다.'

    items = []
    for i, item in enumerate(raw_list):
        if not isinstance(item, dict):
            return None, f'attachments[{i}] 형식이 올바르지 않습니다.'
        fn = (item.get('filename') or item.get('name') or f'file{i + 1}').strip() or f'file{i + 1}'
        mime = (item.get('mime_type') or item.get('mime') or '').strip()
        b64 = item.get('data') or item.get('base64') or ''
        if isinstance(b64, str) and ',' in b64 and b64.startswith('data:'):
            b64 = b64.split(',', 1)[1]
        try:
            raw = base64.standard_b64decode(b64) if b64 else b''
        except Exception:
            return None, f'첨부 `{fn}`의 base64 데이터가 올바르지 않습니다.'
        if len(raw) > CHAT_ATTACHMENT_MAX_BYTES:
            return None, f'첨부 `{fn}`이(가) 너무 큽니다 (최대 {CHAT_ATTACHMENT_MAX_BYTES // (1024 * 1024)}MB).'
        items.append({'filename': fn, 'mime_type': mime, 'raw_bytes': raw})

    return items, None


def merge_last_user_message_with_attachments(messages, attachment_items):
    """
    messages 복사본에서 마지막 user 메시지에 첨부를 반영.
    attachment_items가 비어 있으면 messages 그대로 (복사).
    """
    messages = [dict(m) if isinstance(m, dict) else m for m in (messages or [])]
    if not attachment_items:
        return messages

    last_idx = None
    for i in range(len(messages) - 1, -1, -1):
        m = messages[i]
        if isinstance(m, dict) and m.get('role') == 'user' and not m.get('tool_call_id'):
            last_idx = i
            break
    if last_idx is None:
        messages.append(_build_user_message_with_attachments('', attachment_items))
        return messages

    um = dict(messages[last_idx])
    user_text = um.get('content')
    if isinstance(user_text, list):
        text_parts = [p.get('text', '') for p in user_text if isinstance(p, dict) and p.get('type') == 'text']
        user_text = '\n'.join(text_parts)
    elif user_text is None:
        user_text = ''
    else:
        user_text = str(user_text)

    merged = _build_user_message_with_attachments(user_text, attachment_items)
    um['content'] = merged['content']
    um['display_text'] = merged['display_text']
    messages[last_idx] = um
    return messages


# ===== Chat Provider: OpenAI / Local =====

def chat_openai(
    messages_simple,
    space_key=None,
    allowed_tools=None,
    parent_page_context=None,
    comment_target_page_context=None,
    update_target_page_context=None,
    delete_target_page_context=None,
):
    from openai import OpenAI

    fsk = _forced_search_space_key(space_key)

    provider = ai_config['provider']
    if provider == 'local':
        base_url = ai_config['base_url']
        api_key = get_openai_key() or 'no-key'
    else:
        base_url = 'https://api.openai.com/v1'
        api_key = get_openai_key()
    client = OpenAI(api_key=api_key, base_url=base_url)

    oai_messages = [{
        "role": "system",
        "content": get_system_prompt(
            space_key,
            allowed_tools,
            parent_page_context,
            comment_target_page_context,
            update_target_page_context,
            delete_target_page_context,
        ),
    }]
    for m in messages_simple:
        oai_messages.append({"role": m["role"], "content": m["content"]})

    tools_oai = tools_to_openai(allowed_tools)

    gen_kwargs = {
        'temperature': ai_config['temperature'],
        'top_p': ai_config['top_p'],
    }
    # OpenAI 신규 모델(gpt-5.x, o1 등)은 max_completion_tokens만 지원
    if provider == 'openai':
        gen_kwargs['max_completion_tokens'] = ai_config['max_tokens']
    else:
        gen_kwargs['max_tokens'] = ai_config['max_tokens']

    def _resp_model(resp):
        return getattr(resp, 'model', None) or ai_config['model']

    response = client.chat.completions.create(
        model=ai_config['model'],
        messages=oai_messages,
        tools=tools_oai,
        **gen_kwargs,
    )

    confirm_footer = (
        '\n\n---\n**실행하려면** 채팅에 `확인` 또는 `yes`라고 답하거나 **실행** 버튼을 눌러 주세요. '
        '**취소**하려면 `취소` 또는 **취소** 버튼을 눌러 주세요.'
    )

    for _ in range(10):
        choice = response.choices[0]
        if choice.finish_reason != 'tool_calls' and not getattr(choice.message, 'tool_calls', None):
            return (choice.message.content or '', _resp_model(response), None)

        msg = choice.message
        tcs = msg.tool_calls or []
        names = [tc.function.name for tc in tcs]
        if tcs and any(n in TOOLS_REQUIRING_CONFIRMATION for n in names):
            pending_calls = []
            for tc in tcs:
                try:
                    arg_obj = json.loads(tc.function.arguments or '{}')
                except json.JSONDecodeError:
                    arg_obj = {}
                pending_calls.append({'name': tc.function.name, 'arguments': arg_obj})
            preview_core = format_pending_confirmation_markdown(pending_calls, fsk)
            prefix = (msg.content or '').strip()
            if prefix:
                full_text = f'{prefix}\n\n{preview_core}{confirm_footer}'
            else:
                full_text = preview_core + confirm_footer
            return (full_text, _resp_model(response), pending_calls)

        am = {"role": "assistant", "content": msg.content}
        if tcs:
            am["tool_calls"] = [
                {"id": tc.id, "type": "function", "function": {"name": tc.function.name, "arguments": tc.function.arguments}}
                for tc in tcs
            ]
        oai_messages.append(am)

        for tc in tcs:
            try:
                args = json.loads(tc.function.arguments or '{}')
            except json.JSONDecodeError:
                args = {}
            result = execute_tool(tc.function.name, args, forced_scope_space_key=fsk)
            oai_messages.append({"role": "tool", "tool_call_id": tc.id, "content": result})

        response = client.chat.completions.create(
            model=ai_config['model'],
            messages=oai_messages,
            tools=tools_oai,
            **gen_kwargs,
        )

    return (response.choices[0].message.content or '', _resp_model(response), None)



# ===== Routes =====

def ai_ready():
    return bool(get_openai_key())


@app.route('/api/login', methods=['POST'])
def api_login():
    body = request.get_json() or {}
    raw_id = (body.get('user_id') or body.get('id') or '').strip()
    cf_tok = (body.get('confluence_token') or '').strip()
    oai = (body.get('openai_api_key') or body.get('openai_token') or '').strip()
    if not raw_id or not cf_tok or not oai:
        return jsonify({'error': 'ID, Confluence 토큰, OpenAI API 키를 모두 입력하세요.'}), 400
    if not CONFLUENCE_URL:
        return jsonify({'error': '서버에 CONFLUENCE_URL이 설정되어 있지 않습니다.'}), 503
    username = _normalize_confluence_login_id(raw_id)
    try:
        base, mode, greeting_name = probe_confluence_login(username, cf_tok, raw_id)
    except Exception as e:
        return jsonify({'error': f'Confluence: ID·토큰이 맞지 않거나 API에 접근할 수 없습니다. {e}'}), 401
    try:
        openai_models_ok = verify_openai_key_for_login(oai)
    except Exception as e:
        return jsonify({'error': f'OpenAI: {e}'}), 401

    session.clear()
    session['cf_token'] = cf_tok
    session['cf_auth_mode'] = mode
    session['cf_api_base'] = base
    session['cf_username'] = username if mode == 'basic' else ''
    session['cf_login_id'] = raw_id
    session['cf_greeting_name'] = greeting_name
    session['openai_api_key'] = oai
    return jsonify({
        'success': True,
        'confluenceUser': username if mode == 'basic' else raw_id,
        'authMode': mode,
        'openaiModelsVerified': openai_models_ok,
        'greetingName': greeting_name,
    })


@app.route('/api/logout', methods=['POST'])
def api_logout():
    session.clear()
    return jsonify({'success': True})


def _session_greeting_name():
    raw = (session.get('cf_greeting_name') or '').strip()
    if raw:
        tok = _greeting_token_from_display_name(raw)
        if tok:
            return tok
        return raw
    rid = (session.get('cf_login_id') or '').strip()
    if rid:
        return rid.split('@', 1)[0] if '@' in rid else rid
    return ''


@app.route('/api/status')
def status():
    return jsonify({
        'configured': is_configured(),
        'confluenceUrl': CONFLUENCE_URL or None,
        'aiReady': ai_ready(),
        'loggedIn': bool(session.get('cf_token')),
        'loginId': session.get('cf_login_id'),
        'greetingName': _session_greeting_name(),
        'provider': ai_config['provider'],
    })


@app.route('/api/models')
def list_models():
    provider = request.args.get('provider', ai_config['provider'])

    if provider == 'openai':
        api_key = get_openai_key()
        if not api_key:
            return jsonify({'models': [], 'message': 'API key를 설정하세요'})
        try:
            resp = http_req.get(
                'https://api.openai.com/v1/models',
                headers={'Authorization': f'Bearer {api_key}'},
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            api_ids = {m.get('id', '') for m in data.get('data', [])}
            models = [{'id': mid, 'name': mid} for mid in ALLOWED_OPENAI_MODELS if mid in api_ids]
            if not models:
                models = [{'id': mid, 'name': mid} for mid in ALLOWED_OPENAI_MODELS]
            return jsonify({'models': models})
        except Exception as e:
            return jsonify({'models': [], 'error': str(e)})

    elif provider == 'local':
        base_url = request.args.get('base_url') or ai_config['base_url']
        if not base_url:
            return jsonify({'models': [], 'message': 'Base URL을 먼저 입력하세요'})
        try:
            headers = {}
            api_key = get_openai_key()
            if api_key:
                headers['Authorization'] = f'Bearer {api_key}'
            resp = http_req.get(f"{base_url.rstrip('/')}/models", headers=headers, timeout=10)
            resp.raise_for_status()
            data = resp.json()
            models = sorted([{'id': m['id'], 'name': m['id']} for m in data.get('data', [])], key=lambda x: x['id'])
            return jsonify({'models': models})
        except Exception as e:
            return jsonify({'models': [], 'error': str(e)})

    return jsonify({'models': []})


GEN_PARAMS = ['temperature', 'top_p', 'repetition_penalty', 'max_tokens']

@app.route('/api/config', methods=['GET'])
def get_config():
    cfg = {
        'provider': ai_config['provider'],
        'base_url': ai_config['base_url'],
        'model': ai_config['model'],
        'has_key': bool(get_openai_key()),
    }
    for k in GEN_PARAMS:
        cfg[k] = ai_config[k]
    return jsonify(cfg)


@app.route('/api/config', methods=['POST'])
def set_config():
    body = request.get_json() or {}
    if 'model' in body and body.get('model'):
        m = str(body['model']).strip()
        if m in ALLOWED_OPENAI_MODELS_SET:
            ai_config['model'] = m
    ai_config['provider'] = 'openai'
    cfg = {
        'success': True,
        'provider': 'openai',
        'base_url': '',
        'model': ai_config['model'],
        'has_key': bool(get_openai_key()),
    }
    for k in GEN_PARAMS:
        cfg[k] = ai_config[k]
    return jsonify(cfg)


def fetch_all_confluence_spaces():
    """Load every space (global + personal) with pagination; dedupe by key."""
    merged = {}

    def pull(sp_type):
        start = 0
        limit = 100
        while True:
            path = f'/space?limit={limit}&start={start}'
            if sp_type:
                path += f'&type={sp_type}'
            try:
                data = confluence_api('GET', path)
            except Exception:
                break
            results = data.get('results', [])
            for s in results:
                k = s.get('key')
                if k:
                    st = (s.get('type') or sp_type or 'global')
                    if isinstance(st, str):
                        st = st.lower()
                    merged[k] = {'key': k, 'name': s.get('name') or k, 'type': st}
            if len(results) < limit:
                break
            start += limit

    pull('global')
    pull('personal')
    if not merged:
        pull(None)
    return sorted(merged.values(), key=lambda x: (x.get('name') or x['key']).lower())


@app.route('/api/spaces')
def get_spaces_route():
    try:
        return jsonify({'spaces': fetch_all_confluence_spaces()})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


def _parse_ui_page_context(raw):
    if not raw or not isinstance(raw, dict):
        return None
    ctx = {
        'id': str(raw.get('id') or '').strip(),
        'title': str(raw.get('title') or '').strip(),
        'space_key': str(raw.get('space_key') or '').strip(),
    }
    return ctx if ctx['id'] else None


@app.route('/api/pages/search')
def api_pages_search():
    if not is_configured():
        return jsonify({'error': 'Confluence가 설정되지 않았습니다.', 'pages': []}), 503
    q = (request.args.get('q') or '').strip()
    if len(q) < 2:
        return jsonify({'pages': []})
    space_key = (request.args.get('space_key') or '').strip() or None
    try:
        limit = int(request.args.get('limit', 15))
    except ValueError:
        limit = 15
    limit = max(1, min(limit, 25))
    try:
        raw = tool_search(q, limit=limit, force_space_key=space_key)
        data = json.loads(raw)
    except Exception as e:
        return jsonify({'error': str(e), 'pages': []}), 500
    if isinstance(data, dict) and data.get('error'):
        return jsonify({'error': data['error'], 'pages': []}), 500
    pages = []
    for item in data if isinstance(data, list) else []:
        sid = item.get('id')
        if not sid:
            continue
        sp = item.get('space') or {}
        pages.append({
            'id': str(sid),
            'title': item.get('title') or '',
            'url': item.get('url') or '',
            'space_key': sp.get('key') or '',
            'space_name': sp.get('name') or '',
        })
    return jsonify({'pages': pages})


@app.route('/api/pages/lookup')
def api_pages_lookup():
    if not is_configured():
        return jsonify({'error': 'Confluence가 설정되지 않았습니다.'}), 503
    q = (request.args.get('q') or '').strip()
    page_id = extract_confluence_page_id(q)
    if not page_id:
        return jsonify({'error': '페이지 ID나 Confluence URL을 입력하세요. (예: .../wiki/spaces/X/pages/12345/... )'}), 400
    space_key = (request.args.get('space_key') or '').strip() or None
    allowed = _parse_force_space_keys(space_key)
    try:
        data = confluence_api('GET', f'/content/{page_id}?expand=body.storage,version,space,ancestors,metadata.labels')
    except Exception as e:
        return jsonify({'error': str(e)}), 400
    if allowed is not None:
        pk = _content_space_key(data)
        if pk not in allowed:
            return jsonify({
                'error': (
                    f'이 페이지는 스페이스 "{pk}"에 있습니다. 채팅에서 허용된 스페이스 범위를 선택하거나 전체 스페이스로 두세요.'
                ),
            }), 400
    page = _build_page_result(data, include_content=False)
    return jsonify({'page': page})


def _summarize_pending_execution_results(result_chunks, model_id):
    """도구 실행 직후 원시 결과를 한국어 요약으로 바꿉니다. 실패 시 원문을 돌려줍니다."""
    from openai import OpenAI
    if not get_openai_key() or not result_chunks:
        return '\n\n---\n\n'.join(result_chunks)
    base_url = 'https://api.openai.com/v1'
    client = OpenAI(api_key=get_openai_key(), base_url=base_url)
    mid = model_id if model_id in ALLOWED_OPENAI_MODELS_SET else 'gpt-5.4-mini'
    user_blob = '\n\n---\n\n'.join(result_chunks)
    if len(user_blob) > 12000:
        user_blob = user_blob[:12000] + '\n\n…_(일부 생략)_'
    try:
        gen_kwargs = {'temperature': 0.2, 'max_completion_tokens': 2048}
        r = client.chat.completions.create(
            model=mid,
            messages=[
                {
                    'role': 'system',
                    'content': 'Confluence API 도구 실행 결과를 사용자에게 한국어로 간단히 요약합니다. 성공/실패, 페이지 제목·URL·ID, 오류 메시지 등 핵심만 짧게 bullet으로 알려 주세요.',
                },
                {'role': 'user', 'content': f'다음은 순서대로 실행한 도구 결과입니다:\n\n{user_blob}'},
            ],
            **gen_kwargs,
        )
        t = (r.choices[0].message.content or '').strip()
        return t if t else user_blob
    except Exception:
        return user_blob


@app.route('/api/chat', methods=['POST'])
def chat():
    body = request.get_json()
    messages = body.get('messages', [])
    if not messages:
        return jsonify({'error': 'No messages provided'}), 400

    attachment_items, att_err = parse_chat_attachments_from_request(body or {})
    if att_err:
        return jsonify({'error': att_err}), 400
    messages = merge_last_user_message_with_attachments(messages, attachment_items or [])
    if attachment_items:
        print(
            f'[chat] merged {len(attachment_items)} attachment(s) into last user message',
            flush=True,
        )

    space_key = body.get('space_key', '')
    enabled_tools = normalize_enabled_tools(body.get('enabled_tools'))
    req_model = (body.get('model') or '').strip()
    parent_page_context = _parse_ui_page_context(body.get('parent_page_context'))
    comment_target_page_context = _parse_ui_page_context(body.get('comment_target_page_context'))
    update_target_page_context = _parse_ui_page_context(body.get('update_target_page_context'))
    delete_target_page_context = _parse_ui_page_context(body.get('delete_target_page_context'))

    if not get_openai_key():
        return jsonify({'error': 'OpenAI API 키가 없습니다. 로그아웃 후 다시 로그인하세요.'}), 400

    ai_config['provider'] = 'openai'
    if ai_config['model'] not in ALLOWED_OPENAI_MODELS_SET:
        ai_config['model'] = 'gpt-5.4-mini'
    original_model = ai_config['model']
    if req_model:
        if req_model not in ALLOWED_OPENAI_MODELS_SET:
            return jsonify({
                'error': f'허용된 모델만 사용할 수 있습니다: {", ".join(ALLOWED_OPENAI_MODELS)}',
            }), 400
        ai_config['model'] = req_model

    try:
        messages_for_llm = normalize_messages_for_chat_api(messages)
        final_text, used_model, pending = chat_openai(
            messages_for_llm,
            space_key=space_key or None,
            allowed_tools=enabled_tools,
            parent_page_context=parent_page_context,
            comment_target_page_context=comment_target_page_context,
            update_target_page_context=update_target_page_context,
            delete_target_page_context=delete_target_page_context,
        )
        final_text = re.sub(r'<think>.*?</think>\s*', '', final_text, flags=re.DOTALL)
        print(f"\n===== LLM FINAL RESPONSE (first 3000 chars) =====", flush=True)
        print(final_text[:3000], flush=True)
        print(f"===== END LLM RESPONSE =====\n", flush=True)
        updated = [_strip_llm_only_keys(m) for m in messages] + [{'role': 'assistant', 'content': final_text}]
        out = {'response': final_text, 'messages': updated, 'used_model': used_model}
        if pending:
            out['needs_confirmation'] = True
            out['pending_tool_calls'] = pending
        return jsonify(out)
    except Exception as e:
        traceback.print_exc()
        err = str(e)
        if any(kw in err.lower() for kw in ['auth', '401', 'api key', 'permission']):
            return jsonify({'error': f'Authentication error: {err}'}), 401
        return jsonify({'error': err}), 500
    finally:
        ai_config['model'] = original_model


@app.route('/api/chat/confirm', methods=['POST'])
def chat_confirm_pending():
    """대기 중인 쓰기 도구만 순서대로 실행하고, 결과를 한국어로 요약해 반환합니다."""
    if not is_configured():
        return jsonify({'error': 'Confluence가 설정되지 않았습니다.'}), 503
    if not get_openai_key():
        return jsonify({'error': 'OpenAI API 키가 없습니다. 로그아웃 후 다시 로그인하세요.'}), 400

    body = request.get_json() or {}
    calls = body.get('pending_tool_calls')
    if not calls or not isinstance(calls, list):
        return jsonify({'error': 'pending_tool_calls 배열이 필요합니다.'}), 400

    space_key = body.get('space_key', '')
    enabled_tools = normalize_enabled_tools(body.get('enabled_tools'))
    req_model = (body.get('model') or '').strip()
    if req_model and req_model not in ALLOWED_OPENAI_MODELS_SET:
        return jsonify({
            'error': f'허용된 모델만 사용할 수 있습니다: {", ".join(ALLOWED_OPENAI_MODELS)}',
        }), 400
    summarize_model = req_model or ai_config['model']
    if summarize_model not in ALLOWED_OPENAI_MODELS_SET:
        summarize_model = 'gpt-5.4-mini'

    fsk = _forced_search_space_key(space_key or None)
    normalized = []
    for i, item in enumerate(calls):
        if not isinstance(item, dict):
            return jsonify({'error': f'pending_tool_calls[{i}] 형식이 올바르지 않습니다.'}), 400
        name = (item.get('name') or '').strip()
        if name not in TOOLS_REQUIRING_CONFIRMATION:
            return jsonify({'error': f'이 경로에서 실행할 수 없는 도구입니다: {name or "(이름 없음)"}'}), 400
        if name not in TOOL_MAP:
            return jsonify({'error': f'알 수 없는 도구: {name}'}), 400
        args = item.get('arguments')
        if not isinstance(args, dict):
            args = {}
        if name not in enabled_tools:
            return jsonify({'error': f'현재 채팅 작업 설정에서 허용되지 않는 도구입니다: {name}'}), 400
        normalized.append({'name': name, 'arguments': args})

    try:
        result_chunks = []
        for c in normalized:
            result_str = execute_tool(c['name'], c['arguments'], forced_scope_space_key=fsk)
            result_chunks.append(f"**{c['name']}**\n{result_str}")

        summary = _summarize_pending_execution_results(result_chunks, summarize_model)
        return jsonify({'response': summary, 'used_model': summarize_model})
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/')
@app.route('/<path:path>')
def serve_app(path=''):
    return send_from_directory('public', 'index.html')



if __name__ == '__main__':
    print('Starting Confluence Chat...')
    if not os.getenv('FLASK_SECRET_KEY'):
        print('WARNING: Set FLASK_SECRET_KEY in .env so session cookies persist across restarts.')
    print(f'OpenAI model (default): {ai_config["model"]}')
    print('NOTE: Confluence credentials are expected from browser login (session).')

    host = '127.0.0.1'
    port = int(os.getenv('PORT', 3000))
    try:
        from waitress import serve
        print(f'Server: Waitress threads=8, http://{host}:{port}')
        serve(app, host=host, port=port, threads=8)
    except ImportError:
        print('TIP: pip install waitress  — reduces noisy "Socket is not connected" (Errno 57) on disconnect.')
        print(f'Server: Flask http://{host}:{port}')
        app.run(host=host, port=port, debug=False, threaded=True)
